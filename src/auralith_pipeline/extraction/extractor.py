"""Content extraction module for various file formats."""

from dataclasses import dataclass
from typing import Optional, List, Dict, Any
from pathlib import Path
import logging
import mimetypes

logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """Extracted content from a file."""
    
    text: str
    metadata: Dict[str, Any]
    tables: List[List[List[str]]] = None
    images: List[bytes] = None
    
    def __post_init__(self):
        if self.tables is None:
            self.tables = []
        if self.images is None:
            self.images = []


class ContentExtractor:
    """Extract content from various file formats."""
    
    def __init__(self):
        self._extractors = {
            ".txt": self._extract_text,
            ".md": self._extract_text,
            ".json": self._extract_json,
            ".jsonl": self._extract_jsonl,
            ".csv": self._extract_csv,
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
            ".html": self._extract_html,
            ".xml": self._extract_xml,
        }
    
    def extract(self, file_path: str) -> Optional[ExtractedContent]:
        """Extract content from file."""
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        if suffix in self._extractors:
            try:
                return self._extractors[suffix](path)
            except Exception as e:
                logger.error(f"Failed to extract {file_path}: {e}")
                return None
        else:
            logger.warning(f"Unsupported file type: {suffix}")
            return None
    
    def _extract_text(self, path: Path) -> ExtractedContent:
        """Extract plain text."""
        text = path.read_text(encoding="utf-8", errors="ignore")
        return ExtractedContent(
            text=text,
            metadata={"type": "text", "size": path.stat().st_size},
        )
    
    def _extract_json(self, path: Path) -> ExtractedContent:
        """Extract JSON content."""
        import json
        
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        # Convert to text
        if isinstance(data, dict):
            text = data.get("text", data.get("content", json.dumps(data, indent=2)))
        elif isinstance(data, list):
            texts = [item.get("text", str(item)) if isinstance(item, dict) else str(item) for item in data]
            text = "\n\n".join(texts)
        else:
            text = str(data)
        
        return ExtractedContent(
            text=text,
            metadata={"type": "json", "size": path.stat().st_size},
        )
    
    def _extract_jsonl(self, path: Path) -> ExtractedContent:
        """Extract JSONL content."""
        import json
        
        texts = []
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                try:
                    item = json.loads(line)
                    if isinstance(item, dict):
                        texts.append(item.get("text", str(item)))
                    else:
                        texts.append(str(item))
                except json.JSONDecodeError:
                    continue
        
        return ExtractedContent(
            text="\n\n".join(texts),
            metadata={"type": "jsonl", "num_records": len(texts)},
        )
    
    def _extract_csv(self, path: Path) -> ExtractedContent:
        """Extract CSV content."""
        import csv
        
        rows = []
        with open(path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        
        # Convert to text representation
        text_rows = [", ".join(row) for row in rows]
        
        return ExtractedContent(
            text="\n".join(text_rows),
            metadata={"type": "csv", "num_rows": len(rows)},
            tables=[rows],
        )
    
    def _extract_pdf(self, path: Path) -> ExtractedContent:
        """Extract PDF content using pdfplumber."""
        try:
            import pdfplumber
        except ImportError:
            logger.error("pdfplumber not installed. Install with: pip install pdfplumber")
            return ExtractedContent(text="", metadata={"error": "pdfplumber not installed"})
        
        text_parts = []
        tables = []
        
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                # Extract text
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
                
                # Extract tables
                page_tables = page.extract_tables()
                for table in page_tables:
                    tables.append(table)
        
        return ExtractedContent(
            text="\n\n".join(text_parts),
            metadata={
                "type": "pdf",
                "num_pages": len(text_parts),
                "num_tables": len(tables),
            },
            tables=tables,
        )
    
    def _extract_docx(self, path: Path) -> ExtractedContent:
        """Extract DOCX content."""
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return ExtractedContent(text="", metadata={"error": "python-docx not installed"})
        
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                table_data.append([cell.text for cell in row.cells])
            tables.append(table_data)
        
        return ExtractedContent(
            text="\n\n".join(paragraphs),
            metadata={
                "type": "docx",
                "num_paragraphs": len(paragraphs),
                "num_tables": len(tables),
            },
            tables=tables,
        )
    
    def _extract_pptx(self, path: Path) -> ExtractedContent:
        """Extract PPTX content."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return ExtractedContent(text="", metadata={"error": "python-pptx not installed"})
        
        prs = Presentation(path)
        texts = []
        
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append("\n".join(slide_texts))
        
        return ExtractedContent(
            text="\n\n---\n\n".join(texts),  # Separate slides
            metadata={"type": "pptx", "num_slides": len(texts)},
        )
    
    def _extract_xlsx(self, path: Path) -> ExtractedContent:
        """Extract XLSX content."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return ExtractedContent(text="", metadata={"error": "openpyxl not installed"})
        
        wb = load_workbook(path, read_only=True)
        texts = []
        tables = []
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            sheet_data = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                sheet_data.append(row_data)
            
            if sheet_data:
                tables.append(sheet_data)
                # Convert to text
                text_rows = [", ".join(row) for row in sheet_data]
                texts.append(f"Sheet: {sheet}\n" + "\n".join(text_rows))
        
        return ExtractedContent(
            text="\n\n".join(texts),
            metadata={"type": "xlsx", "num_sheets": len(texts)},
            tables=tables,
        )
    
    def _extract_html(self, path: Path) -> ExtractedContent:
        """Extract HTML content."""
        try:
            from html.parser import HTMLParser
            
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.texts = []
                    self.skip_tags = {"script", "style", "head", "meta", "link"}
                    self._skip = False
                
                def handle_starttag(self, tag, attrs):
                    if tag in self.skip_tags:
                        self._skip = True
                
                def handle_endtag(self, tag):
                    if tag in self.skip_tags:
                        self._skip = False
                
                def handle_data(self, data):
                    if not self._skip and data.strip():
                        self.texts.append(data.strip())
            
            html_content = path.read_text(encoding="utf-8", errors="ignore")
            parser = TextExtractor()
            parser.feed(html_content)
            
            return ExtractedContent(
                text="\n".join(parser.texts),
                metadata={"type": "html", "size": len(html_content)},
            )
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            return ExtractedContent(text="", metadata={"error": str(e)})
    
    def _extract_xml(self, path: Path) -> ExtractedContent:
        """Extract XML content."""
        import xml.etree.ElementTree as ET
        
        tree = ET.parse(path)
        root = tree.getroot()
        
        def get_text(element):
            texts = []
            if element.text and element.text.strip():
                texts.append(element.text.strip())
            for child in element:
                texts.extend(get_text(child))
            return texts
        
        texts = get_text(root)
        
        return ExtractedContent(
            text="\n".join(texts),
            metadata={"type": "xml", "root_tag": root.tag},
        )

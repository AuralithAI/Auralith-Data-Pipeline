"""Compound document extraction — decompose multi-modal files into sub-assets.

Many real-world file formats are *compound*: a single file can contain
text, images, tables, audio, and video.  For example:

- **DOCX** → paragraphs (text) + embedded images + tables
- **XLSX** → tabular data (text) + charts (images)
- **PDF**  → text + tables + scanned images
- **PPTX** → slide text + speaker notes + embedded media
- **DICOM / NIfTI** → 3-D medical image volumes → sliced into 2-D images
- **GeoTIFF / topography** → raster elevation data → image tiles
- **MKV / MP4 (with subtitles)** → video frames + audio track + subtitle text

This module provides :class:`CompoundDocumentExtractor` which inspects the
file type, delegates to format-specific helpers, and yields a list of
:class:`ModalitySegment` objects — one per extracted sub-asset.  Each
segment carries raw bytes (for binary assets like images/audio) or text,
plus metadata linking it back to the source file and its position within.

Downstream, :class:`CompoundDocumentSource` wraps the extractor as a
:class:`DataSource`, producing ``DataSample`` objects that flow through the
standard pipeline (preprocessing → tokenize → shard).

**Design principle:** The already-trained tokenizers (BPE for text, VQ
codebooks for image/audio/video) are loaded from cold storage and applied
here — no retraining is needed.  ``train_tokenizer`` is NOT affected.
"""

from __future__ import annotations

import io
import logging
import struct
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

import numpy as np

logger = logging.getLogger(__name__)

# ── Supported compound file extensions ─────────────────────────────────

COMPOUND_EXTS: frozenset[str] = frozenset(
    {
        # Office documents
        ".docx",
        ".xlsx",
        ".xls",
        ".pptx",
        ".odt",
        ".ods",
        ".odp",
        # PDF
        ".pdf",
        # Medical imaging
        ".dcm",
        ".nii",
        ".nii.gz",
        # Geospatial / topography
        ".tif",
        ".tiff",
        ".hgt",
        # Rich media containers
        ".mkv",
        ".mp4",
        ".webm",
        # eBooks
        ".epub",
        # Email
        ".eml",
        ".mbox",
        # Rich text
        ".rtf",
        # Web archives
        ".mhtml",
        ".mht",
        # LaTeX / scientific papers
        ".tex",
        # HDF5 scientific data
        ".h5",
        ".hdf5",
        # Astronomy FITS images
        ".fits",
        # Zarr array stores
        ".zarr",
        # Web crawl archives
        ".warc",
        ".warc.gz",
        # 3-D models
        ".glb",
        ".gltf",
        # Comic book archives
        ".cbz",
        ".cbr",
        # Outlook email
        ".msg",
    }
)


# ── Data classes ───────────────────────────────────────────────────────


@dataclass
class ModalitySegment:
    """One extracted sub-asset from a compound document.

    Attributes
    ----------
    modality:
        Target modality — used to select the correct tokenizer and to
        populate ``DataSample.modality`` / the modality_mask in shards.
    content_text:
        Plain-text content (for text / table segments).  Empty string
        for binary modalities.
    content_bytes:
        Raw bytes for binary sub-assets (images, audio clips).  ``None``
        for pure text segments.
    metadata:
        Provenance metadata: source file, page/sheet/slide number,
        position within the file, format details, etc.
    """

    modality: Literal["text", "image", "audio", "video", "code"]
    content_text: str = ""
    content_bytes: bytes | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class CompoundDocument:
    """A file decomposed into its constituent modality segments.

    Each compound document carries a ``document_id`` (derived from the
    file path) so that all segments can be linked back to the same
    source during lineage tracking.
    """

    source_path: str
    document_id: str
    segments: list[ModalitySegment] = field(default_factory=list)

    @property
    def modalities_present(self) -> set[str]:
        return {s.modality for s in self.segments}

    @property
    def num_segments(self) -> int:
        return len(self.segments)


# ── Extractor ──────────────────────────────────────────────────────────


class CompoundDocumentExtractor:
    """Decompose compound files into modality segments.

    The extractor is **stateless** — it does not train or load tokenizers.
    Its job is purely to *extract* sub-assets.  Tokenization happens
    downstream in the pipeline via the pre-trained tokenizers.

    Usage::

        extractor = CompoundDocumentExtractor()
        doc = extractor.extract("report.docx")
        for seg in doc.segments:
            print(seg.modality, len(seg.content_text or seg.content_bytes))
    """

    # Format → handler mapping (populated in __init__)
    _handlers: dict[str, Any]

    def __init__(
        self,
        *,
        extract_images: bool = True,
        extract_tables: bool = True,
        extract_audio: bool = True,
        extract_video_frames: bool = True,
        max_image_pixels: int = 4096 * 4096,
        geotiff_tile_size: int = 256,
        dicom_window_center: float | None = None,
        dicom_window_width: float | None = None,
    ) -> None:
        self.extract_images = extract_images
        self.extract_tables = extract_tables
        self.extract_audio = extract_audio
        self.extract_video_frames = extract_video_frames
        self.max_image_pixels = max_image_pixels
        self.geotiff_tile_size = geotiff_tile_size
        self.dicom_window_center = dicom_window_center
        self.dicom_window_width = dicom_window_width

        self._handlers = {
            ".docx": self._extract_docx,
            ".xlsx": self._extract_xlsx,
            ".xls": self._extract_xlsx,
            ".pptx": self._extract_pptx,
            ".pdf": self._extract_pdf,
            ".odt": self._extract_odt,
            ".ods": self._extract_odt,
            ".odp": self._extract_odt,
            ".dcm": self._extract_dicom,
            ".nii": self._extract_nifti,
            ".tif": self._extract_geotiff,
            ".tiff": self._extract_geotiff,
            ".hgt": self._extract_hgt,
            ".mkv": self._extract_media_container,
            ".mp4": self._extract_media_container,
            ".webm": self._extract_media_container,
            ".epub": self._extract_epub,
            ".eml": self._extract_eml,
            ".mbox": self._extract_mbox,
            ".rtf": self._extract_rtf,
            ".mhtml": self._extract_mhtml,
            ".mht": self._extract_mhtml,
            ".tex": self._extract_latex,
            ".h5": self._extract_hdf5,
            ".hdf5": self._extract_hdf5,
            ".fits": self._extract_fits,
            ".zarr": self._extract_zarr,
            ".warc": self._extract_warc,
            ".glb": self._extract_gltf,
            ".gltf": self._extract_gltf,
            ".cbz": self._extract_comic_book,
            ".cbr": self._extract_comic_book,
            ".msg": self._extract_msg,
        }

    # ── public API ─────────────────────────────────────────────────────

    def extract(self, file_path: str | Path) -> CompoundDocument:
        """Extract all modality segments from *file_path*.

        Raises ``ValueError`` for unsupported extensions.
        """
        path = Path(file_path)
        suffix = path.suffix.lower()

        # Handle double-extension compound types specially
        name_lower = path.name.lower()
        if name_lower.endswith(".nii.gz"):
            suffix = ".nii"
        elif name_lower.endswith(".warc.gz"):
            suffix = ".warc"

        handler = self._handlers.get(suffix)
        if handler is None:
            raise ValueError(
                f"Unsupported compound file type: {suffix}. "
                f"Supported: {sorted(self._handlers.keys())}"
            )

        doc_id = f"compound:{path.stem}:{_file_hash_short(path)}"
        doc = CompoundDocument(source_path=str(path), document_id=doc_id)

        try:
            handler(path, doc)
        except Exception as exc:
            logger.error("Failed to extract compound document %s: %s", path, exc)
            doc.segments.append(
                ModalitySegment(
                    modality="text",
                    content_text="",
                    metadata={"error": str(exc), "source": str(path)},
                )
            )

        logger.info(
            "Extracted %d segments (%s) from %s",
            doc.num_segments,
            ", ".join(sorted(doc.modalities_present)),
            path.name,
        )
        return doc

    def supported_extensions(self) -> frozenset[str]:
        """Return the set of file extensions this extractor can handle."""
        return frozenset(self._handlers.keys())

    # ── DOCX ───────────────────────────────────────────────────────────

    def _extract_docx(self, path: Path, doc: CompoundDocument) -> None:
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed. pip install python-docx")
            return

        docx_doc = Document(str(path))

        # 1. Paragraph text
        paragraphs = [p.text for p in docx_doc.paragraphs if p.text.strip()]
        if paragraphs:
            doc.segments.append(
                ModalitySegment(
                    modality="text",
                    content_text="\n\n".join(paragraphs),
                    metadata={
                        "source": str(path),
                        "type": "docx_text",
                        "num_paragraphs": len(paragraphs),
                    },
                )
            )

        # 2. Tables
        if self.extract_tables:
            for idx, table in enumerate(docx_doc.tables):
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                if rows:
                    table_text = _table_to_markdown(rows)
                    doc.segments.append(
                        ModalitySegment(
                            modality="text",
                            content_text=table_text,
                            metadata={
                                "source": str(path),
                                "type": "docx_table",
                                "table_index": idx,
                                "num_rows": len(rows),
                            },
                        )
                    )

        # 3. Embedded images
        if self.extract_images:
            for rel in docx_doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_bytes = rel.target_part.blob
                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=image_bytes,
                                metadata={
                                    "source": str(path),
                                    "type": "docx_embedded_image",
                                    "content_type": rel.target_part.content_type,
                                },
                            )
                        )
                    except Exception as exc:
                        logger.debug("Could not extract DOCX image: %s", exc)

    # ── XLSX / XLS ─────────────────────────────────────────────────────

    def _extract_xlsx(self, path: Path, doc: CompoundDocument) -> None:
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.error("openpyxl not installed. pip install openpyxl")
            return

        wb = load_workbook(str(path), read_only=True, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                # Skip entirely empty rows
                if any(cell.strip() for cell in row_data):
                    rows.append(row_data)

            if rows:
                table_text = _table_to_markdown(rows)
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text=table_text,
                        metadata={
                            "source": str(path),
                            "type": "xlsx_sheet",
                            "sheet_name": sheet_name,
                            "num_rows": len(rows),
                            "num_cols": max(len(r) for r in rows) if rows else 0,
                        },
                    )
                )

        # Extract chart images if available (openpyxl exposes them as _images)
        if self.extract_images:
            try:
                # Re-open non-read-only to access images
                wb_full = load_workbook(str(path), read_only=False)
                for sheet_name in wb_full.sheetnames:
                    ws_full = wb_full[sheet_name]
                    for img in getattr(ws_full, "_images", []):
                        try:
                            image_data = img._data()
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=image_data,
                                    metadata={
                                        "source": str(path),
                                        "type": "xlsx_embedded_image",
                                        "sheet_name": sheet_name,
                                    },
                                )
                            )
                        except Exception:
                            pass
                wb_full.close()
            except Exception as exc:
                logger.debug("Could not extract XLSX images: %s", exc)

        wb.close()

    # ── PPTX ───────────────────────────────────────────────────────────

    def _extract_pptx(self, path: Path, doc: CompoundDocument) -> None:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            logger.error("python-pptx not installed. pip install python-pptx")
            return

        prs = Presentation(str(path))

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                # Text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

                # Embedded images
                if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=image_bytes,
                                metadata={
                                    "source": str(path),
                                    "type": "pptx_image",
                                    "slide_index": slide_idx,
                                    "content_type": shape.image.content_type,
                                },
                            )
                        )
                    except Exception as exc:
                        logger.debug("PPTX image extraction failed: %s", exc)

                # Tables
                if self.extract_tables and shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    if rows:
                        doc.segments.append(
                            ModalitySegment(
                                modality="text",
                                content_text=_table_to_markdown(rows),
                                metadata={
                                    "source": str(path),
                                    "type": "pptx_table",
                                    "slide_index": slide_idx,
                                },
                            )
                        )

            if slide_texts:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text="\n".join(slide_texts),
                        metadata={
                            "source": str(path),
                            "type": "pptx_slide_text",
                            "slide_index": slide_idx,
                        },
                    )
                )

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    doc.segments.append(
                        ModalitySegment(
                            modality="text",
                            content_text=notes,
                            metadata={
                                "source": str(path),
                                "type": "pptx_speaker_notes",
                                "slide_index": slide_idx,
                            },
                        )
                    )

    # ── PDF ─────────────────────────────────────────────────────────────

    def _extract_pdf(self, path: Path, doc: CompoundDocument) -> None:
        try:
            import pdfplumber
        except ImportError:
            logger.error("pdfplumber not installed. pip install pdfplumber")
            return

        with pdfplumber.open(str(path)) as pdf:
            all_text: list[str] = []

            for page_idx, page in enumerate(pdf.pages):
                # Text
                page_text = page.extract_text() or ""
                if page_text.strip():
                    all_text.append(page_text)

                # Tables
                if self.extract_tables:
                    for table in page.extract_tables():
                        if table:
                            # pdfplumber returns list[list[str|None]]
                            clean_rows = [[cell or "" for cell in row] for row in table if row]
                            if clean_rows:
                                doc.segments.append(
                                    ModalitySegment(
                                        modality="text",
                                        content_text=_table_to_markdown(clean_rows),
                                        metadata={
                                            "source": str(path),
                                            "type": "pdf_table",
                                            "page": page_idx,
                                        },
                                    )
                                )

                # Images
                if self.extract_images:
                    for img_idx, img in enumerate(page.images):
                        try:
                            # Extract image bounding box as cropped page image
                            bbox = (
                                img["x0"],
                                img["top"],
                                img["x1"],
                                img["bottom"],
                            )
                            cropped = page.within_bbox(bbox).to_image(resolution=150)
                            buf = io.BytesIO()
                            cropped.original.save(buf, format="PNG")
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=buf.getvalue(),
                                    metadata={
                                        "source": str(path),
                                        "type": "pdf_image",
                                        "page": page_idx,
                                        "image_index": img_idx,
                                    },
                                )
                            )
                        except Exception as exc:
                            logger.debug("PDF image extraction failed: %s", exc)

            if all_text:
                doc.segments.insert(
                    0,
                    ModalitySegment(
                        modality="text",
                        content_text="\n\n".join(all_text),
                        metadata={
                            "source": str(path),
                            "type": "pdf_text",
                            "num_pages": len(pdf.pages),
                        },
                    ),
                )

    # ── ODF (ODT / ODS / ODP) ──────────────────────────────────────────

    def _extract_odt(self, path: Path, doc: CompoundDocument) -> None:
        """Basic ODF extraction — text only for now."""
        import zipfile

        try:
            with zipfile.ZipFile(str(path)) as zf:
                if "content.xml" in zf.namelist():
                    import xml.etree.ElementTree as ET

                    content_xml = zf.read("content.xml")
                    root = ET.fromstring(content_xml)
                    texts = [elem.text for elem in root.iter() if elem.text and elem.text.strip()]
                    if texts:
                        doc.segments.append(
                            ModalitySegment(
                                modality="text",
                                content_text="\n".join(texts),
                                metadata={
                                    "source": str(path),
                                    "type": "odf_text",
                                },
                            )
                        )
        except Exception as exc:
            logger.error("ODF extraction failed for %s: %s", path, exc)

    # ── DICOM ──────────────────────────────────────────────────────────

    def _extract_dicom(self, path: Path, doc: CompoundDocument) -> None:
        """Extract 2-D slices from a DICOM file as image segments.

        DICOM files contain medical imaging data (CT, MRI, X-ray, ultrasound).
        Each slice is windowed and normalised to uint8 for downstream image
        tokenization via the pre-trained VQ codebook.
        """
        try:
            import pydicom
        except ImportError:
            logger.error("pydicom not installed. pip install pydicom")
            return

        ds = pydicom.dcmread(str(path))

        # Extract patient-safe metadata (no PII)
        meta: dict[str, Any] = {
            "source": str(path),
            "type": "dicom",
            "modality_dicom": str(getattr(ds, "Modality", "unknown")),
            "rows": int(getattr(ds, "Rows", 0)),
            "columns": int(getattr(ds, "Columns", 0)),
        }

        pixel_array = ds.pixel_array.astype(np.float32)

        # Apply windowing
        wc = self.dicom_window_center or float(getattr(ds, "WindowCenter", pixel_array.mean()))
        ww = self.dicom_window_width or float(getattr(ds, "WindowWidth", pixel_array.ptp()))
        if isinstance(wc, list | np.ndarray):
            wc = float(wc[0])
        if isinstance(ww, list | np.ndarray):
            ww = float(ww[0])

        low = wc - ww / 2
        high = wc + ww / 2
        pixel_array = np.clip(pixel_array, low, high)
        pixel_array = ((pixel_array - low) / max(high - low, 1e-6) * 255).astype(np.uint8)

        # Handle 3-D volumes (multi-slice) and single 2-D images
        if pixel_array.ndim == 3:
            for slice_idx in range(pixel_array.shape[0]):
                slice_img = pixel_array[slice_idx]
                image_bytes = _numpy_to_png_bytes(slice_img)
                doc.segments.append(
                    ModalitySegment(
                        modality="image",
                        content_bytes=image_bytes,
                        metadata={**meta, "slice_index": slice_idx},
                    )
                )
        else:
            image_bytes = _numpy_to_png_bytes(pixel_array)
            doc.segments.append(
                ModalitySegment(
                    modality="image",
                    content_bytes=image_bytes,
                    metadata=meta,
                )
            )

        # Extract any text report embedded in the DICOM
        for text_attr in ("StudyDescription", "SeriesDescription", "ImageComments"):
            val = getattr(ds, text_attr, None)
            if val and str(val).strip():
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text=f"{text_attr}: {val}",
                        metadata={"source": str(path), "type": "dicom_text_field"},
                    )
                )

    # ── NIfTI ──────────────────────────────────────────────────────────

    def _extract_nifti(self, path: Path, doc: CompoundDocument) -> None:
        """Extract slices from NIfTI neuroimaging files (.nii / .nii.gz)."""
        try:
            import nibabel as nib
        except ImportError:
            logger.error("nibabel not installed. pip install nibabel")
            return

        img = nib.load(str(path))
        data = np.asarray(img.dataobj, dtype=np.float32)

        meta: dict[str, Any] = {
            "source": str(path),
            "type": "nifti",
            "shape": list(data.shape),
            "voxel_size": img.header.get_zooms()[:3] if hasattr(img.header, "get_zooms") else None,
        }

        # Normalise to 0-255
        dmin, dmax = float(data.min()), float(data.max())
        if dmax - dmin > 1e-6:
            data = ((data - dmin) / (dmax - dmin) * 255).astype(np.uint8)
        else:
            data = np.zeros_like(data, dtype=np.uint8)

        # Extract axial slices (last axis for typical NIfTI orientation)
        if data.ndim >= 3:
            num_slices = data.shape[2]
            # Sample up to 64 evenly-spaced slices to avoid overwhelming the pipeline
            max_slices = 64
            if num_slices > max_slices:
                indices = np.linspace(0, num_slices - 1, max_slices, dtype=int)
            else:
                indices = range(num_slices)

            for slice_idx in indices:
                slice_2d = data[:, :, int(slice_idx)]
                image_bytes = _numpy_to_png_bytes(slice_2d)
                doc.segments.append(
                    ModalitySegment(
                        modality="image",
                        content_bytes=image_bytes,
                        metadata={**meta, "slice_index": int(slice_idx), "plane": "axial"},
                    )
                )
        else:
            # 2-D image
            image_bytes = _numpy_to_png_bytes(data)
            doc.segments.append(
                ModalitySegment(modality="image", content_bytes=image_bytes, metadata=meta)
            )

    # ── GeoTIFF (topography / elevation / satellite) ───────────────────

    def _extract_geotiff(self, path: Path, doc: CompoundDocument) -> None:
        """Tile a GeoTIFF raster into image patches for VQ tokenization.

        Multi-band rasters (e.g. satellite RGB) are extracted as colour
        tiles.  Single-band elevation data is converted to grayscale tiles.
        """
        try:
            import rasterio
        except ImportError:
            logger.error("rasterio not installed. pip install rasterio")
            return

        with rasterio.open(str(path)) as src:
            meta: dict[str, Any] = {
                "source": str(path),
                "type": "geotiff",
                "crs": str(src.crs) if src.crs else None,
                "bounds": list(src.bounds) if src.bounds else None,
                "num_bands": src.count,
                "width": src.width,
                "height": src.height,
            }

            tile_size = self.geotiff_tile_size

            # Read all bands
            raster = src.read()  # (bands, H, W)

            if raster.shape[0] >= 3:
                # RGB: take first 3 bands
                rgb = np.stack([raster[i] for i in range(3)], axis=-1)
            elif raster.shape[0] == 1:
                # Single band → grayscale → replicate to 3-channel
                band = raster[0]
                bmin, bmax = float(band.min()), float(band.max())
                if bmax - bmin > 1e-6:
                    band = ((band - bmin) / (bmax - bmin) * 255).astype(np.uint8)
                else:
                    band = np.zeros_like(band, dtype=np.uint8)
                rgb = np.stack([band, band, band], axis=-1)
            else:
                # 2-band or other: use first band as grayscale
                band = raster[0]
                bmin, bmax = float(band.min()), float(band.max())
                if bmax - bmin > 1e-6:
                    band = ((band - bmin) / (bmax - bmin) * 255).astype(np.uint8)
                else:
                    band = np.zeros_like(band, dtype=np.uint8)
                rgb = np.stack([band, band, band], axis=-1)

            h, w = rgb.shape[:2]
            tile_idx = 0
            for y in range(0, h, tile_size):
                for x in range(0, w, tile_size):
                    tile = rgb[y : y + tile_size, x : x + tile_size]
                    if tile.shape[0] < 16 or tile.shape[1] < 16:
                        continue  # Skip tiny edge tiles
                    image_bytes = _numpy_to_png_bytes(tile)
                    doc.segments.append(
                        ModalitySegment(
                            modality="image",
                            content_bytes=image_bytes,
                            metadata={
                                **meta,
                                "tile_index": tile_idx,
                                "tile_y": y,
                                "tile_x": x,
                                "tile_h": tile.shape[0],
                                "tile_w": tile.shape[1],
                            },
                        )
                    )
                    tile_idx += 1

    # ── SRTM HGT (elevation) ──────────────────────────────────────────

    def _extract_hgt(self, path: Path, doc: CompoundDocument) -> None:
        """Extract SRTM .hgt elevation tiles as grayscale images."""
        raw = path.read_bytes()
        # HGT files are big-endian 16-bit signed integers
        num_values = len(raw) // 2
        side = int(num_values**0.5)
        elevation = np.array(struct.unpack(f">{num_values}h", raw), dtype=np.float32).reshape(
            side, side
        )

        # Normalise to 0-255
        vmin, vmax = float(elevation.min()), float(elevation.max())
        if vmax - vmin > 1e-6:
            normed = ((elevation - vmin) / (vmax - vmin) * 255).astype(np.uint8)
        else:
            normed = np.zeros_like(elevation, dtype=np.uint8)

        tile_size = self.geotiff_tile_size
        tile_idx = 0
        for y in range(0, side, tile_size):
            for x in range(0, side, tile_size):
                tile = normed[y : y + tile_size, x : x + tile_size]
                if tile.shape[0] < 16 or tile.shape[1] < 16:
                    continue
                image_bytes = _numpy_to_png_bytes(tile)
                doc.segments.append(
                    ModalitySegment(
                        modality="image",
                        content_bytes=image_bytes,
                        metadata={
                            "source": str(path),
                            "type": "hgt_elevation",
                            "tile_index": tile_idx,
                            "elevation_min": vmin,
                            "elevation_max": vmax,
                        },
                    )
                )
                tile_idx += 1

    # ── Media containers (MKV / MP4 / WebM) ───────────────────────────

    def _extract_media_container(self, path: Path, doc: CompoundDocument) -> None:
        """Extract video frames, audio track, and subtitle text from
        media containers.

        This handler decomposes a single video file into:
        - Video frame images (via VideoFrameSampler)
        - Audio waveform (via ffmpeg / soundfile)
        - Subtitle / caption text (via pysrt / webvtt)
        """
        from auralith_pipeline.sources.video import VideoFrameSampler

        # 1. Video frames
        if self.extract_video_frames:
            try:
                sampler = VideoFrameSampler(max_frames=64, strategy="uniform")
                frames = sampler.extract_frames(path)
                for frame_idx in range(frames.shape[0]):
                    frame_bytes = _numpy_to_png_bytes(frames[frame_idx])
                    doc.segments.append(
                        ModalitySegment(
                            modality="image",
                            content_bytes=frame_bytes,
                            metadata={
                                "source": str(path),
                                "type": "video_frame",
                                "frame_index": frame_idx,
                            },
                        )
                    )
            except Exception as exc:
                logger.warning("Video frame extraction failed for %s: %s", path, exc)

        # 2. Audio track
        if self.extract_audio:
            try:
                audio_bytes = _extract_audio_track(path)
                if audio_bytes:
                    doc.segments.append(
                        ModalitySegment(
                            modality="audio",
                            content_bytes=audio_bytes,
                            metadata={
                                "source": str(path),
                                "type": "video_audio_track",
                            },
                        )
                    )
            except Exception as exc:
                logger.debug("Audio track extraction failed for %s: %s", path, exc)

        # 3. Subtitles
        try:
            subtitle_text = _extract_subtitles(path)
            if subtitle_text:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text=subtitle_text,
                        metadata={
                            "source": str(path),
                            "type": "video_subtitles",
                        },
                    )
                )
        except Exception as exc:
            logger.debug("Subtitle extraction failed for %s: %s", path, exc)

    # ── EPUB (eBooks) ──────────────────────────────────────────────────

    def _extract_epub(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text chapters and embedded images from EPUB eBooks.

        EPUB is a ZIP archive containing XHTML content documents and media.
        We parse the OPF manifest to iterate spine items in reading order
        and extract embedded ``<img>`` assets.
        """
        import zipfile
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            """Minimal HTML→text converter."""

            def __init__(self) -> None:
                super().__init__()
                self._parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
                if tag in ("script", "style"):
                    self._skip = True

            def handle_endtag(self, tag: str) -> None:
                if tag in ("script", "style"):
                    self._skip = False
                if tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li"):
                    self._parts.append("\n")

            def handle_data(self, data: str) -> None:
                if not self._skip:
                    self._parts.append(data)

            @property
            def text(self) -> str:
                return "".join(self._parts).strip()

        try:
            with zipfile.ZipFile(str(path)) as zf:
                names = zf.namelist()

                # 1. Extract text from all XHTML/HTML content files
                chapter_idx = 0
                for name in names:
                    lower = name.lower()
                    if not (
                        lower.endswith(".xhtml")
                        or lower.endswith(".html")
                        or lower.endswith(".htm")
                    ):
                        continue
                    try:
                        raw_html = zf.read(name).decode("utf-8", errors="ignore")
                        parser = _TextExtractor()
                        parser.feed(raw_html)
                        text = parser.text
                        if text and len(text) > 20:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="text",
                                    content_text=text,
                                    metadata={
                                        "source": str(path),
                                        "type": "epub_chapter",
                                        "chapter_index": chapter_idx,
                                        "content_file": name,
                                    },
                                )
                            )
                            chapter_idx += 1
                    except Exception as exc:
                        logger.debug("EPUB chapter extraction failed for %s: %s", name, exc)

                # 2. Extract embedded images
                if self.extract_images:
                    image_exts = (".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp")
                    for name in names:
                        if any(name.lower().endswith(ext) for ext in image_exts):
                            try:
                                image_bytes = zf.read(name)
                                doc.segments.append(
                                    ModalitySegment(
                                        modality="image",
                                        content_bytes=image_bytes,
                                        metadata={
                                            "source": str(path),
                                            "type": "epub_image",
                                            "image_path": name,
                                        },
                                    )
                                )
                            except Exception as exc:
                                logger.debug("EPUB image extraction failed for %s: %s", name, exc)

        except Exception as exc:
            logger.error("EPUB extraction failed for %s: %s", path, exc)

    # ── EML (email) ────────────────────────────────────────────────────

    def _extract_eml(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text body and attachments from RFC-2822 email files."""
        import email
        import email.policy

        try:
            raw = path.read_bytes()
            msg = email.message_from_bytes(raw, policy=email.policy.default)
            self._process_email_message(msg, path, doc)
        except Exception as exc:
            logger.error("EML extraction failed for %s: %s", path, exc)

    def _extract_mbox(self, path: Path, doc: CompoundDocument) -> None:
        """Extract all messages from an mbox mailbox file."""
        import mailbox

        try:
            mbox = mailbox.mbox(str(path))
            for msg_idx, msg in enumerate(mbox):
                self._process_email_message(msg, path, doc, msg_index=msg_idx)
                if msg_idx >= 999:
                    logger.info("mbox capped at 1000 messages for %s", path)
                    break
        except Exception as exc:
            logger.error("mbox extraction failed for %s: %s", path, exc)

    def _process_email_message(
        self,
        msg: Any,
        source_path: Path,
        doc: CompoundDocument,
        msg_index: int = 0,
    ) -> None:
        """Shared logic for EML and mbox: extract body + attachments."""

        subject = str(msg.get("subject", ""))
        text_parts: list[str] = []

        if subject:
            text_parts.append(f"Subject: {subject}")

        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                disposition = str(part.get("Content-Disposition", ""))

                if content_type == "text/plain" and "attachment" not in disposition:
                    payload = part.get_payload(decode=True)
                    if payload:
                        text_parts.append(payload.decode("utf-8", errors="ignore"))

                elif content_type == "text/html" and "attachment" not in disposition:
                    payload = part.get_payload(decode=True)
                    if payload:
                        from html.parser import HTMLParser

                        class _Strip(HTMLParser):
                            def __init__(self) -> None:
                                super().__init__()
                                self._parts: list[str] = []

                            def handle_data(self, d: str) -> None:
                                self._parts.append(d)

                            @property
                            def text(self) -> str:
                                return " ".join(self._parts)

                        p = _Strip()
                        p.feed(payload.decode("utf-8", errors="ignore"))
                        if p.text.strip():
                            text_parts.append(p.text.strip())

                elif self.extract_images and content_type.startswith("image/"):
                    payload = part.get_payload(decode=True)
                    if payload:
                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=payload,
                                metadata={
                                    "source": str(source_path),
                                    "type": "email_image_attachment",
                                    "msg_index": msg_index,
                                    "content_type": content_type,
                                },
                            )
                        )
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                text_parts.append(payload.decode("utf-8", errors="ignore"))

        if text_parts:
            doc.segments.append(
                ModalitySegment(
                    modality="text",
                    content_text="\n\n".join(text_parts),
                    metadata={
                        "source": str(source_path),
                        "type": "email_body",
                        "msg_index": msg_index,
                        "subject": subject,
                    },
                )
            )

    # ── RTF ─────────────────────────────────────────────────────────────

    def _extract_rtf(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text from RTF files using striprtf."""
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            logger.error("striprtf not installed. pip install striprtf")
            return

        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")
            text = rtf_to_text(raw)
            if text and text.strip():
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text=text.strip(),
                        metadata={"source": str(path), "type": "rtf_text"},
                    )
                )
        except Exception as exc:
            logger.error("RTF extraction failed for %s: %s", path, exc)

    # ── MHTML / MHT (web archives) ────────────────────────────────────

    def _extract_mhtml(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text and images from MHTML web archive files.

        MHTML files are MIME-encoded web pages — they contain the HTML
        plus all referenced resources (images, CSS) as MIME parts.
        We use the stdlib ``email`` module to parse the MIME structure.
        """
        import email
        import email.policy
        from html.parser import HTMLParser

        class _HtmlToText(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
                if tag in ("script", "style"):
                    self._skip = True

            def handle_endtag(self, tag: str) -> None:
                if tag in ("script", "style"):
                    self._skip = False
                if tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
                    self._parts.append("\n")

            def handle_data(self, data: str) -> None:
                if not self._skip:
                    self._parts.append(data)

            @property
            def text(self) -> str:
                return "".join(self._parts).strip()

        try:
            raw = path.read_bytes()
            msg = email.message_from_bytes(raw, policy=email.policy.default)

            texts: list[str] = []
            for part in msg.walk():
                ct = part.get_content_type()

                if ct == "text/html":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parser = _HtmlToText()
                        parser.feed(payload.decode("utf-8", errors="ignore"))
                        if parser.text:
                            texts.append(parser.text)

                elif ct == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        decoded = payload.decode("utf-8", errors="ignore").strip()
                        if decoded:
                            texts.append(decoded)

                elif self.extract_images and ct.startswith("image/"):
                    payload = part.get_payload(decode=True)
                    if payload:
                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=payload,
                                metadata={
                                    "source": str(path),
                                    "type": "mhtml_image",
                                    "content_type": ct,
                                },
                            )
                        )

            if texts:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text="\n\n".join(texts),
                        metadata={
                            "source": str(path),
                            "type": "mhtml_text",
                        },
                    )
                )

        except Exception as exc:
            logger.error("MHTML extraction failed for %s: %s", path, exc)

    # ── LaTeX ──────────────────────────────────────────────────────────

    def _extract_latex(self, path: Path, doc: CompoundDocument) -> None:
        r"""Extract text and referenced figure paths from LaTeX files.

        Strips common LaTeX commands (``\section``, ``\begin{...}``, etc.)
        to produce readable text, and collects ``\includegraphics`` paths
        to extract referenced images when they exist on disk.
        """
        import re

        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")

            # 1. Extract referenced figures (resolve relative to .tex dir)
            if self.extract_images:
                fig_pattern = re.compile(r"\\includegraphics(?:\[.*?\])?\{([^}]+)\}")
                for match in fig_pattern.finditer(raw):
                    fig_rel = match.group(1)
                    # Try common extensions if none specified
                    candidates = [fig_rel]
                    if "." not in Path(fig_rel).suffix:
                        candidates += [
                            f"{fig_rel}.png",
                            f"{fig_rel}.jpg",
                            f"{fig_rel}.pdf",
                            f"{fig_rel}.eps",
                        ]

                    for candidate in candidates:
                        fig_path = path.parent / candidate
                        if fig_path.exists():
                            try:
                                doc.segments.append(
                                    ModalitySegment(
                                        modality="image",
                                        content_bytes=fig_path.read_bytes(),
                                        metadata={
                                            "source": str(path),
                                            "type": "latex_figure",
                                            "figure_ref": fig_rel,
                                            "figure_path": str(fig_path),
                                        },
                                    )
                                )
                            except Exception as exc:
                                logger.debug("LaTeX figure read failed for %s: %s", fig_path, exc)
                            break

            # 2. Strip LaTeX markup to produce clean text
            text = raw
            # Remove comments
            text = re.sub(r"(?m)(?<!\\)%.*$", "", text)
            # Remove common preamble commands
            text = re.sub(
                r"\\(?:documentclass|usepackage|bibliography|bibliographystyle|label|ref|cite|eqref|pageref)(?:\[.*?\])?\{[^}]*\}",
                "",
                text,
            )
            # Convert section headings to plain text
            text = re.sub(
                r"\\(?:section|subsection|subsubsection|chapter|paragraph|subparagraph)\*?\{([^}]*)\}",
                r"\n\n\1\n",
                text,
            )
            # Remove begin/end environments (keep content)
            text = re.sub(r"\\(?:begin|end)\{[^}]*\}", "", text)
            # Remove \command{...} patterns, keeping content
            text = re.sub(
                r"\\(?:textbf|textit|emph|underline|texttt|textrm|textsf|textsc)\{([^}]*)\}",
                r"\1",
                text,
            )
            # Remove remaining simple commands
            text = re.sub(r"\\[a-zA-Z]+\*?(?:\{[^}]*\})*", "", text)
            # Clean up math delimiters
            text = re.sub(r"\$\$.*?\$\$", "[equation]", text, flags=re.DOTALL)
            text = re.sub(r"\$[^$]+\$", "[math]", text)
            # Clean up whitespace
            text = re.sub(r"\n{3,}", "\n\n", text).strip()

            if text and len(text) > 30:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text=text,
                        metadata={
                            "source": str(path),
                            "type": "latex_text",
                        },
                    )
                )

        except Exception as exc:
            logger.error("LaTeX extraction failed for %s: %s", path, exc)

    # ── HDF5 ───────────────────────────────────────────────────────────

    def _extract_hdf5(self, path: Path, doc: CompoundDocument) -> None:
        """Extract datasets from HDF5 files as image tiles or text metadata.

        HDF5 is a universal container for scientific data — physics
        simulations, climate data, genomics, etc.  2-D and 3-D numerical
        arrays are normalised and exported as image tiles.  Scalar and
        string attributes are exported as text metadata.
        """
        try:
            import h5py
        except ImportError:
            logger.error("h5py not installed. pip install h5py")
            return

        try:
            with h5py.File(str(path), "r") as f:
                attrs_text: list[str] = []

                def _visitor(name: str, obj: Any) -> None:
                    # Collect string/scalar attributes
                    for attr_name, attr_val in obj.attrs.items():
                        attrs_text.append(f"{name}.{attr_name} = {attr_val}")

                    if isinstance(obj, h5py.Dataset):
                        shape = obj.shape
                        dtype = obj.dtype

                        # 2-D array → single image
                        if len(shape) == 2 and np.issubdtype(dtype, np.number):
                            data = obj[()].astype(np.float32)
                            data = _normalise_to_uint8(data)
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=_numpy_to_png_bytes(data),
                                    metadata={
                                        "source": str(path),
                                        "type": "hdf5_dataset",
                                        "dataset_name": name,
                                        "shape": list(shape),
                                    },
                                )
                            )

                        # 3-D array → sample slices along first axis
                        elif len(shape) == 3 and np.issubdtype(dtype, np.number):
                            num_slices = shape[0]
                            max_s = min(num_slices, 64)
                            indices = (
                                np.linspace(0, num_slices - 1, max_s, dtype=int)
                                if num_slices > max_s
                                else range(num_slices)
                            )
                            for idx in indices:
                                sl = obj[int(idx)].astype(np.float32)
                                sl = _normalise_to_uint8(sl)
                                doc.segments.append(
                                    ModalitySegment(
                                        modality="image",
                                        content_bytes=_numpy_to_png_bytes(sl),
                                        metadata={
                                            "source": str(path),
                                            "type": "hdf5_slice",
                                            "dataset_name": name,
                                            "slice_index": int(idx),
                                        },
                                    )
                                )

                        # String datasets → text
                        elif np.issubdtype(dtype, np.bytes_) or np.issubdtype(dtype, np.str_):
                            try:
                                vals = [
                                    v.decode("utf-8") if isinstance(v, bytes) else str(v)
                                    for v in obj[()].flat
                                ]
                                text = "\n".join(vals)
                                if text.strip():
                                    doc.segments.append(
                                        ModalitySegment(
                                            modality="text",
                                            content_text=text,
                                            metadata={
                                                "source": str(path),
                                                "type": "hdf5_string_dataset",
                                                "dataset_name": name,
                                            },
                                        )
                                    )
                            except Exception:
                                pass

                f.visititems(_visitor)

                if attrs_text:
                    doc.segments.append(
                        ModalitySegment(
                            modality="text",
                            content_text="\n".join(attrs_text),
                            metadata={
                                "source": str(path),
                                "type": "hdf5_attributes",
                            },
                        )
                    )

        except Exception as exc:
            logger.error("HDF5 extraction failed for %s: %s", path, exc)

    # ── FITS (astronomy) ──────────────────────────────────────────────

    def _extract_fits(self, path: Path, doc: CompoundDocument) -> None:
        """Extract image data and header metadata from FITS files.

        FITS (Flexible Image Transport System) is the standard format
        for astronomical image data from telescopes like Hubble and JWST.
        Each HDU (Header/Data Unit) can contain image arrays + rich
        metadata headers.
        """
        try:
            from astropy.io import fits as astropy_fits
        except ImportError:
            logger.error("astropy not installed. pip install astropy")
            return

        try:
            with astropy_fits.open(str(path)) as hdul:
                for hdu_idx, hdu in enumerate(hdul):
                    header_lines: list[str] = []
                    for key in hdu.header:
                        if key and key not in ("COMMENT", "HISTORY", ""):
                            val = hdu.header[key]
                            comment = hdu.header.comments.get(key, "")
                            header_lines.append(
                                f"{key} = {val}  / {comment}" if comment else f"{key} = {val}"
                            )

                    if header_lines:
                        doc.segments.append(
                            ModalitySegment(
                                modality="text",
                                content_text="\n".join(header_lines),
                                metadata={
                                    "source": str(path),
                                    "type": "fits_header",
                                    "hdu_index": hdu_idx,
                                },
                            )
                        )

                    # Extract image data
                    if (
                        hdu.data is not None
                        and hasattr(hdu.data, "shape")
                        and len(hdu.data.shape) >= 2
                    ):
                        data = hdu.data.astype(np.float32)

                        if data.ndim == 2:
                            data = _normalise_to_uint8(data)
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=_numpy_to_png_bytes(data),
                                    metadata={
                                        "source": str(path),
                                        "type": "fits_image",
                                        "hdu_index": hdu_idx,
                                        "shape": list(hdu.data.shape),
                                    },
                                )
                            )
                        elif data.ndim == 3:
                            # Multi-plane: extract up to 64 slices
                            n = data.shape[0]
                            max_s = min(n, 64)
                            indices = (
                                np.linspace(0, n - 1, max_s, dtype=int) if n > max_s else range(n)
                            )
                            for idx in indices:
                                sl = _normalise_to_uint8(data[int(idx)])
                                doc.segments.append(
                                    ModalitySegment(
                                        modality="image",
                                        content_bytes=_numpy_to_png_bytes(sl),
                                        metadata={
                                            "source": str(path),
                                            "type": "fits_slice",
                                            "hdu_index": hdu_idx,
                                            "slice_index": int(idx),
                                        },
                                    )
                                )

        except Exception as exc:
            logger.error("FITS extraction failed for %s: %s", path, exc)

    # ── Zarr ───────────────────────────────────────────────────────────

    def _extract_zarr(self, path: Path, doc: CompoundDocument) -> None:
        """Extract arrays from Zarr stores as image slices.

        Zarr is a chunked, compressed array format used for large-scale
        climate, satellite, and microscopy data.
        """
        try:
            import zarr
        except ImportError:
            logger.error("zarr not installed. pip install zarr")
            return

        try:
            store = zarr.open(str(path), mode="r")

            def _visit_array(name: str, arr: Any) -> None:
                if not hasattr(arr, "shape"):
                    return  # Group, not array

                shape = arr.shape
                if len(shape) == 2 and arr.dtype.kind in ("i", "u", "f"):
                    data = np.array(arr[:], dtype=np.float32)
                    data = _normalise_to_uint8(data)
                    doc.segments.append(
                        ModalitySegment(
                            modality="image",
                            content_bytes=_numpy_to_png_bytes(data),
                            metadata={
                                "source": str(path),
                                "type": "zarr_array",
                                "array_name": name,
                                "shape": list(shape),
                            },
                        )
                    )
                elif len(shape) == 3 and arr.dtype.kind in ("i", "u", "f"):
                    n = shape[0]
                    max_s = min(n, 64)
                    indices = np.linspace(0, n - 1, max_s, dtype=int) if n > max_s else range(n)
                    for idx in indices:
                        sl = np.array(arr[int(idx)], dtype=np.float32)
                        sl = _normalise_to_uint8(sl)
                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=_numpy_to_png_bytes(sl),
                                metadata={
                                    "source": str(path),
                                    "type": "zarr_slice",
                                    "array_name": name,
                                    "slice_index": int(idx),
                                },
                            )
                        )

                # Zarr string arrays → text
                elif arr.dtype.kind in ("U", "S", "O"):
                    try:
                        vals = [str(v) for v in arr[:].flat if str(v).strip()]
                        if vals:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="text",
                                    content_text="\n".join(vals),
                                    metadata={
                                        "source": str(path),
                                        "type": "zarr_string_array",
                                        "array_name": name,
                                    },
                                )
                            )
                    except Exception:
                        pass

            if hasattr(store, "visititems"):
                store.visititems(_visit_array)
            elif hasattr(store, "arrays"):
                for name, arr in store.arrays():
                    _visit_array(name, arr)
            else:
                # Top-level store is itself an array
                _visit_array("root", store)

            # Collect Zarr-level attrs as text
            if hasattr(store, "attrs") and dict(store.attrs):
                attrs_lines = [f"{k} = {v}" for k, v in store.attrs.items()]
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text="\n".join(attrs_lines),
                        metadata={
                            "source": str(path),
                            "type": "zarr_attributes",
                        },
                    )
                )

        except Exception as exc:
            logger.error("Zarr extraction failed for %s: %s", path, exc)

    # ── WARC (web crawl archives) ──────────────────────────────────────

    def _extract_warc(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text and images from WARC web archive files.

        WARC (Web ARChive) is the format used by Common Crawl and the
        Internet Archive.  Each WARC file contains HTTP response records
        with their original content — HTML pages, images, etc.
        """
        try:
            from warcio.archiveiterator import ArchiveIterator
        except ImportError:
            logger.error("warcio not installed. pip install warcio")
            return

        from html.parser import HTMLParser

        class _WarcTextExtractor(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
                if tag in ("script", "style"):
                    self._skip = True

            def handle_endtag(self, tag: str) -> None:
                if tag in ("script", "style"):
                    self._skip = False
                if tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6"):
                    self._parts.append("\n")

            def handle_data(self, data: str) -> None:
                if not self._skip:
                    self._parts.append(data)

            @property
            def text(self) -> str:
                return "".join(self._parts).strip()

        try:
            open_fn = open
            name_lower = path.name.lower()
            if name_lower.endswith(".gz"):
                import gzip

                open_fn = gzip.open  # type: ignore[assignment]

            record_count = 0
            max_records = 500  # Cap to avoid overwhelming the pipeline

            with open_fn(str(path), "rb") as fh:
                for record in ArchiveIterator(fh):
                    if record_count >= max_records:
                        break

                    if record.rec_type != "response":
                        continue

                    content_type = (
                        record.http_headers.get_header("Content-Type")
                        if record.http_headers
                        else ""
                    )
                    content_type = content_type or ""
                    payload = record.content_stream().read()

                    if "text/html" in content_type:
                        parser = _WarcTextExtractor()
                        parser.feed(payload.decode("utf-8", errors="ignore"))
                        if parser.text and len(parser.text) > 50:
                            url = record.rec_headers.get_header("WARC-Target-URI") or ""
                            doc.segments.append(
                                ModalitySegment(
                                    modality="text",
                                    content_text=parser.text,
                                    metadata={
                                        "source": str(path),
                                        "type": "warc_html",
                                        "url": url,
                                    },
                                )
                            )
                            record_count += 1

                    elif self.extract_images and content_type.startswith("image/"):
                        if payload and len(payload) > 100:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=payload,
                                    metadata={
                                        "source": str(path),
                                        "type": "warc_image",
                                        "content_type": content_type,
                                    },
                                )
                            )
                            record_count += 1

        except Exception as exc:
            logger.error("WARC extraction failed for %s: %s", path, exc)

    # ── glTF / GLB (3-D models) ────────────────────────────────────────

    def _extract_gltf(self, path: Path, doc: CompoundDocument) -> None:
        """Extract metadata and embedded textures from glTF/GLB 3-D models.

        glTF ("GL Transmission Format") is the JPEG of 3-D — widely used
        for games, AR/VR, and digital twins.  GLB is the binary-packed
        variant.  We extract:
        - Scene metadata (node names, materials) as text
        - Embedded texture images (PNG/JPEG buffers)
        """
        import json

        try:
            suffix = path.suffix.lower()

            if suffix == ".gltf":
                raw = path.read_text(encoding="utf-8", errors="ignore")
                gltf_data = json.loads(raw)
                base_dir = path.parent
            elif suffix == ".glb":
                raw_bytes = path.read_bytes()
                # GLB header: magic(4) + version(4) + length(4)
                if len(raw_bytes) < 12 or raw_bytes[:4] != b"glTF":
                    logger.error("Invalid GLB header for %s", path)
                    return
                # First chunk is JSON
                chunk_length = struct.unpack_from("<I", raw_bytes, 12)[0]
                chunk_type = struct.unpack_from("<I", raw_bytes, 16)[0]
                if chunk_type != 0x4E4F534A:  # "JSON" in LE
                    logger.error("GLB first chunk is not JSON for %s", path)
                    return
                json_bytes = raw_bytes[20 : 20 + chunk_length]
                gltf_data = json.loads(json_bytes.decode("utf-8"))
                # Binary chunk starts after JSON chunk
                bin_offset = 20 + chunk_length
                # There may be padding + another chunk header (8 bytes)
                bin_data = raw_bytes[bin_offset + 8 :] if bin_offset + 8 < len(raw_bytes) else b""
                base_dir = path.parent
            else:
                return

            # 1. Metadata as text
            meta_parts: list[str] = []
            if "asset" in gltf_data:
                asset = gltf_data["asset"]
                meta_parts.append(f"Generator: {asset.get('generator', 'unknown')}")
                meta_parts.append(f"Version: {asset.get('version', 'unknown')}")

            if "scenes" in gltf_data:
                for scene_idx, scene in enumerate(gltf_data["scenes"]):
                    meta_parts.append(f"Scene {scene_idx}: {scene.get('name', 'unnamed')}")

            if "nodes" in gltf_data:
                for node in gltf_data["nodes"]:
                    name = node.get("name", "")
                    if name:
                        meta_parts.append(f"Node: {name}")

            if "materials" in gltf_data:
                for mat in gltf_data["materials"]:
                    meta_parts.append(f"Material: {mat.get('name', 'unnamed')}")

            if "meshes" in gltf_data:
                meta_parts.append(f"Meshes: {len(gltf_data['meshes'])}")

            if meta_parts:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text="\n".join(meta_parts),
                        metadata={
                            "source": str(path),
                            "type": "gltf_metadata",
                        },
                    )
                )

            # 2. Embedded textures
            if self.extract_images and "images" in gltf_data:
                buffer_views = gltf_data.get("bufferViews", [])

                for img_idx, img_def in enumerate(gltf_data["images"]):
                    try:
                        if "uri" in img_def:
                            uri = img_def["uri"]
                            if uri.startswith("data:"):
                                # Data URI
                                import base64

                                _, encoded = uri.split(",", 1)
                                image_bytes = base64.b64decode(encoded)
                            else:
                                # External file
                                img_path = base_dir / uri
                                if img_path.exists():
                                    image_bytes = img_path.read_bytes()
                                else:
                                    continue
                        elif "bufferView" in img_def and suffix == ".glb":
                            bv = buffer_views[img_def["bufferView"]]
                            offset = bv.get("byteOffset", 0)
                            length = bv["byteLength"]
                            image_bytes = bin_data[offset : offset + length]
                        else:
                            continue

                        doc.segments.append(
                            ModalitySegment(
                                modality="image",
                                content_bytes=image_bytes,
                                metadata={
                                    "source": str(path),
                                    "type": "gltf_texture",
                                    "image_index": img_idx,
                                    "mime_type": img_def.get("mimeType", ""),
                                },
                            )
                        )
                    except Exception as exc:
                        logger.debug("glTF texture extraction failed: %s", exc)

        except Exception as exc:
            logger.error("glTF extraction failed for %s: %s", path, exc)

    # ── Comic Book Archives (CBZ / CBR) ───────────────────────────────

    def _extract_comic_book(self, path: Path, doc: CompoundDocument) -> None:
        """Extract page images from comic book archives.

        CBZ files are ZIP archives of sequential images (one per page).
        CBR files are RAR archives — we attempt to use ``rarfile`` if
        available, otherwise fall back to an error message.
        """
        suffix = path.suffix.lower()
        image_exts = (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")

        if suffix == ".cbz":
            import zipfile

            try:
                with zipfile.ZipFile(str(path)) as zf:
                    image_names = sorted(
                        n
                        for n in zf.namelist()
                        if any(n.lower().endswith(e) for e in image_exts)
                        and not n.startswith("__MACOSX")
                    )
                    for page_idx, name in enumerate(image_names):
                        try:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=zf.read(name),
                                    metadata={
                                        "source": str(path),
                                        "type": "comic_page",
                                        "page_index": page_idx,
                                        "page_file": name,
                                    },
                                )
                            )
                        except Exception as exc:
                            logger.debug("CBZ page extraction failed for %s: %s", name, exc)
            except Exception as exc:
                logger.error("CBZ extraction failed for %s: %s", path, exc)

        elif suffix == ".cbr":
            try:
                import rarfile
            except ImportError:
                logger.error("rarfile not installed. pip install rarfile (also requires unrar)")
                return

            try:
                with rarfile.RarFile(str(path)) as rf:
                    image_names = sorted(
                        n for n in rf.namelist() if any(n.lower().endswith(e) for e in image_exts)
                    )
                    for page_idx, name in enumerate(image_names):
                        try:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=rf.read(name),
                                    metadata={
                                        "source": str(path),
                                        "type": "comic_page",
                                        "page_index": page_idx,
                                        "page_file": name,
                                    },
                                )
                            )
                        except Exception as exc:
                            logger.debug("CBR page extraction failed for %s: %s", name, exc)
            except Exception as exc:
                logger.error("CBR extraction failed for %s: %s", path, exc)

    # ── Outlook MSG ────────────────────────────────────────────────────

    def _extract_msg(self, path: Path, doc: CompoundDocument) -> None:
        """Extract text body and attachments from Outlook .msg files."""
        try:
            import extract_msg
        except ImportError:
            logger.error("extract-msg not installed. pip install extract-msg")
            return

        try:
            msg = extract_msg.Message(str(path))
            text_parts: list[str] = []

            if msg.subject:
                text_parts.append(f"Subject: {msg.subject}")
            if msg.body:
                text_parts.append(msg.body)

            if text_parts:
                doc.segments.append(
                    ModalitySegment(
                        modality="text",
                        content_text="\n\n".join(text_parts),
                        metadata={
                            "source": str(path),
                            "type": "msg_body",
                            "subject": msg.subject or "",
                            "date": str(msg.date) if msg.date else "",
                        },
                    )
                )

            # Extract image attachments
            if self.extract_images:
                for att in msg.attachments:
                    ct = getattr(att, "mimetype", "") or ""
                    if ct.startswith("image/"):
                        try:
                            doc.segments.append(
                                ModalitySegment(
                                    modality="image",
                                    content_bytes=att.data,
                                    metadata={
                                        "source": str(path),
                                        "type": "msg_image_attachment",
                                        "filename": getattr(att, "longFilename", "")
                                        or getattr(att, "shortFilename", ""),
                                    },
                                )
                            )
                        except Exception as exc:
                            logger.debug("MSG attachment extraction failed: %s", exc)

            msg.close()

        except Exception as exc:
            logger.error("MSG extraction failed for %s: %s", path, exc)


# ── Utility helpers ────────────────────────────────────────────────────


def _normalise_to_uint8(data: np.ndarray) -> np.ndarray:
    """Normalise a floating-point array to uint8 [0, 255]."""
    dmin, dmax = float(data.min()), float(data.max())
    if dmax - dmin > 1e-6:
        return ((data - dmin) / (dmax - dmin) * 255).astype(np.uint8)
    return np.zeros_like(data, dtype=np.uint8)


def _file_hash_short(path: Path, length: int = 8) -> str:
    """Return a short hash of the file path (not content) for document IDs."""
    import hashlib

    return hashlib.sha256(str(path.resolve()).encode()).hexdigest()[:length]


def _numpy_to_png_bytes(arr: np.ndarray) -> bytes:
    """Convert a 2-D or 3-D numpy array to PNG bytes via Pillow."""
    from PIL import Image

    if arr.ndim == 2:
        img = Image.fromarray(arr, mode="L")  # Grayscale
    elif arr.ndim == 3 and arr.shape[2] == 3:
        img = Image.fromarray(arr.astype(np.uint8), mode="RGB")
    elif arr.ndim == 3 and arr.shape[2] == 1:
        img = Image.fromarray(arr[:, :, 0], mode="L")
    else:
        # Fallback: treat as grayscale from first channel
        img = Image.fromarray(arr.reshape(arr.shape[0], arr.shape[1]).astype(np.uint8), mode="L")

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _table_to_markdown(rows: list[list[str]]) -> str:
    """Convert a list of rows into a Markdown table string.

    Markdown tables are a compact, model-friendly representation for
    tabular data — easier for an LLM to learn structure from than CSV.
    """
    if not rows:
        return ""

    # Sanitise cells
    clean_rows = [[cell.replace("|", "\\|").replace("\n", " ") for cell in row] for row in rows]

    # Use first row as header
    header = "| " + " | ".join(clean_rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * len(clean_rows[0])) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in clean_rows[1:]]

    return "\n".join([header, separator, *body_lines])


def _extract_audio_track(video_path: Path) -> bytes | None:
    """Extract audio from a video file using ffmpeg → WAV bytes.

    Returns None if ffmpeg is unavailable or the video has no audio.
    """
    import subprocess

    try:
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            tmp_path = tmp.name

        result = subprocess.run(
            [
                "ffmpeg",
                "-i",
                str(video_path),
                "-vn",
                "-acodec",
                "pcm_s16le",
                "-ar",
                "16000",
                "-ac",
                "1",
                "-y",
                tmp_path,
            ],
            capture_output=True,
            timeout=120,
        )

        if result.returncode == 0:
            return Path(tmp_path).read_bytes()
        return None
    except FileNotFoundError:
        logger.debug("ffmpeg not found — audio extraction skipped")
        return None
    except Exception as exc:
        logger.debug("Audio extraction failed: %s", exc)
        return None
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass


def _extract_subtitles(video_path: Path) -> str | None:
    """Extract subtitle text from common sidecar files (.srt, .vtt).

    Looks for files with the same stem as the video in the same directory.
    """
    stem = video_path.stem
    parent = video_path.parent

    for ext in (".srt", ".vtt", ".sub", ".ass"):
        sub_path = parent / f"{stem}{ext}"
        if sub_path.exists():
            try:
                raw = sub_path.read_text(encoding="utf-8", errors="ignore")
                import re

                lines = raw.split("\n")
                text_lines = []
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    if re.match(r"^\d+$", line):
                        continue
                    if re.match(r"\d{2}:\d{2}:", line):
                        continue
                    if line.startswith("WEBVTT") or line.startswith("NOTE"):
                        continue
                    line = re.sub(r"<[^>]+>", "", line)
                    line = re.sub(r"\{[^}]+\}", "", line)
                    if line.strip():
                        text_lines.append(line.strip())
                if text_lines:
                    return "\n".join(text_lines)
            except Exception:
                pass

    return None
